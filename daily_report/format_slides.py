"""PPTX slide deck formatter for daily-report.

Requires python-pptx: pip install python-pptx
"""

from __future__ import annotations

import re

from pptx import Presentation
from pptx.util import Inches, Pt

from daily_report.report_data import ContentItem, RepoContent, ReportData


def format_slides(report: ReportData, output_path: str, group_by: str = "project") -> None:
    """Render the report as a PPTX slide deck.

    When ``report.consolidated_markdown`` is set, parses the consolidated
    markdown into slides. Otherwise uses structured ``report.content``.

    Args:
        report: Complete report data with content already prepared.
        output_path: File path to write the .pptx file.
        group_by: Grouping mode — "project", "status", or "contribution".

    Raises:
        OSError: If the file cannot be written (permissions, missing directory).
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _add_title_slide(prs, report)

    if report.consolidated_markdown:
        sections = _parse_markdown_sections(report.consolidated_markdown)
        for title, bullets in sections:
            _add_text_slide(prs, title, bullets)
    else:
        for repo in report.content:
            _add_content_slide(prs, repo)

    _add_summary_slide(prs, report)

    prs.save(output_path)


# --- internal helpers (private) ---


def _add_title_slide(prs: Presentation, report: ReportData) -> None:
    """Add the title slide with user and date range."""
    layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Activity Report"
    if report.date_from == report.date_to:
        subtitle_text = f"{report.user}\n{report.date_from}"
    else:
        subtitle_text = f"{report.user}\n{report.date_from} .. {report.date_to}"
    slide.placeholders[1].text = subtitle_text


def _add_content_slide(prs: Presentation, repo: RepoContent) -> None:
    """Add a content slide for a single repository."""
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = repo.repo_name

    tf = slide.placeholders[1].text_frame
    tf.clear()
    first_paragraph = True

    for block in repo.blocks:
        # Block items with inline heading label
        skip_status = {repo.repo_name, block.heading}
        for item in block.items:
            p = tf.paragraphs[0] if first_paragraph else tf.add_paragraph()
            first_paragraph = False
            p.text = f"{block.heading}: {_render_item(item, skip_status)}"
            p.level = 0
            p.runs[0].font.size = Pt(12)


def _add_summary_slide(prs: Presentation, report: ReportData) -> None:
    """Add the summary slide with aggregate metrics."""
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Summary"

    s = report.summary

    if s.ai_summary:
        bullets = [s.ai_summary]
    else:
        themes_str = ", ".join(s.themes) if s.themes else "general development"
        merged_label = "merged" if s.is_range else "merged today"
        bullets = [
            f"Total PRs: {s.total_prs}",
            f"Repositories: {s.repo_count}",
            f"{s.merged_count} {merged_label}",
            f"{s.open_count} still open",
            f"Key themes: {themes_str}",
        ]

    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.runs[0].font.size = Pt(14)


def _render_item(item: ContentItem, skip_status: set[str] | None = None) -> str:
    """Render a ContentItem as plain text for slides."""
    text = item.title

    if item.numbers:
        if len(item.numbers) == 1:
            text += f" #{item.numbers[0]}"
        else:
            refs = ", ".join(f"#{n}" for n in item.numbers)
            text += f" ({refs})"

    if item.author:
        text += f" ({item.author})"

    if item.status and item.status not in (skip_status or set()):
        text += f" -- {item.status}"

    if item.status in ("Open", "Draft") and (item.additions or item.deletions):
        text += f" (+{item.additions}/-{item.deletions})"

    if item.reviewers:
        reviewer_str = ", ".join(item.reviewers)
        text += f" -- reviewer: {reviewer_str}"

    if item.days_waiting:
        text += f" -- {item.days_waiting} days"

    return text


# Strip markdown formatting: bold, italic, links, code
_MD_LINK_RE = re.compile(r"\[([^\]]+)\]\([^)]+\)")
_MD_BOLD_RE = re.compile(r"\*\*([^*]+)\*\*")
_MD_ITALIC_RE = re.compile(r"\*([^*]+)\*")
_MD_CODE_RE = re.compile(r"`([^`]+)`")


def _strip_markdown(text: str) -> str:
    """Remove markdown formatting, returning plain text."""
    text = _MD_LINK_RE.sub(r"\1", text)
    text = _MD_BOLD_RE.sub(r"\1", text)
    text = _MD_ITALIC_RE.sub(r"\1", text)
    text = _MD_CODE_RE.sub(r"\1", text)
    return text.strip()


def _parse_markdown_sections(markdown: str) -> list[tuple[str, list[str]]]:
    """Parse consolidated markdown into (section_title, bullets) pairs.

    Expects H2 (##) sections with ``- `` bullet items. Skips the H1 title
    and the **Summary:** line (handled separately by the summary slide).
    """
    sections: list[tuple[str, list[str]]] = []
    current_title = ""
    current_bullets: list[str] = []

    for line in markdown.splitlines():
        stripped = line.strip()
        if stripped.startswith("# ") and not stripped.startswith("## "):
            continue  # skip H1
        if stripped.startswith("**Summary:**"):
            continue  # skip summary line
        if stripped.startswith("## "):
            if current_title and current_bullets:
                sections.append((current_title, current_bullets))
            current_title = _strip_markdown(stripped[3:].strip())
            current_bullets = []
        elif stripped.startswith("- "):
            bullet_text = _strip_markdown(stripped[2:].strip())
            if bullet_text:
                current_bullets.append(bullet_text)

    if current_title and current_bullets:
        sections.append((current_title, current_bullets))

    return sections


def _add_text_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    """Add a slide with a title and plain-text bullet items."""
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.level = 0
        p.runs[0].font.size = Pt(12)
