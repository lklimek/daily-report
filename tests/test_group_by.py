"""Unit tests for --group-by feature: regroup_content(), and group_by
parameter propagation through Markdown, Slides, and Slack formatters.

Run with: python3 -m pytest tests/test_group_by.py -v
"""

import subprocess
import sys
from pathlib import Path

import pytest
from pptx import Presentation

from daily_report.content import regroup_content
from daily_report.format_markdown import format_markdown
from daily_report.format_slides import format_slides
from daily_report.format_slack import format_slack
from daily_report.report_data import (
    AuthoredPR,
    ContentBlock,
    ContentItem,
    RepoContent,
    ReportData,
    ReviewedPR,
    SummaryStats,
    WaitingPR,
)

PROJECT_ROOT = Path(__file__).resolve().parent.parent


# ---------------------------------------------------------------------------
# Helpers: reusable factories
# ---------------------------------------------------------------------------

def _make_group_by_report() -> ReportData:
    """Create a report with PRs across 2 repos in all three categories."""
    return ReportData(
        user="alice",
        date_from="2026-02-10",
        date_to="2026-02-10",
        authored_prs=[
            AuthoredPR(
                repo="org/alpha", title="Add login", number=10,
                status="Open", additions=50, deletions=10,
                contributed=False, original_author=None,
            ),
            AuthoredPR(
                repo="org/beta", title="Fix crash", number=20,
                status="Merged", additions=0, deletions=0,
                contributed=True, original_author="bob",
            ),
        ],
        reviewed_prs=[
            ReviewedPR(
                repo="org/alpha", title="Update docs", number=11,
                author="charlie", status="Open",
            ),
        ],
        waiting_prs=[
            WaitingPR(
                repo="org/beta", title="Refactor DB", number=21,
                reviewers=["dave", "eve"], created_at="2026-02-08",
                days_waiting=2,
            ),
        ],
        summary=SummaryStats(
            total_prs=4, repo_count=2, merged_count=1,
            open_count=2, themes=["feat", "fix"], is_range=False,
        ),
    )


def _make_regrouped_report(group_by: str) -> ReportData:
    """Create a report and apply regroup_content with the given mode."""
    report = _make_group_by_report()
    report.content = regroup_content(report, group_by)
    return report


def _all_section_texts(blocks):
    """Extract all mrkdwn text values from Slack section blocks."""
    texts = []
    for b in blocks:
        if b.get("type") == "section":
            text_obj = b.get("text", {})
            if text_obj.get("type") == "mrkdwn":
                texts.append(text_obj["text"])
    return texts


# ===========================================================================
# 1. regroup_content() tests
# ===========================================================================


class TestRegroupContribution:
    """Tests for group_by='contribution' mode."""

    @pytest.fixture(autouse=True)
    def _setup(self):
        report = _make_group_by_report()
        self.content = regroup_content(report, "contribution")

    def test_top_level_groups(self):
        names = [rc.repo_name for rc in self.content]
        assert names == ["Authored / Contributed", "Reviewed", "Waiting for Review"]

    def test_authored_group_has_two_repos(self):
        authored = self.content[0]
        headings = [b.heading for b in authored.blocks]
        assert headings == ["org/alpha", "org/beta"]

    def test_reviewed_group_has_one_repo(self):
        reviewed = self.content[1]
        headings = [b.heading for b in reviewed.blocks]
        assert headings == ["org/alpha"]

    def test_waiting_group_has_one_repo(self):
        waiting = self.content[2]
        headings = [b.heading for b in waiting.blocks]
        assert headings == ["org/beta"]

    def test_item_fields_preserved(self):
        # Authored item
        authored_block = self.content[0].blocks[0]  # org/alpha
        item = authored_block.items[0]
        assert item.title == "Add login"
        assert item.status == "Open"
        assert item.additions == 50
        assert item.deletions == 10

        # Reviewed item
        reviewed_block = self.content[1].blocks[0]  # org/alpha
        item = reviewed_block.items[0]
        assert item.author == "charlie"

        # Waiting item
        waiting_block = self.content[2].blocks[0]  # org/beta
        item = waiting_block.items[0]
        assert item.reviewers == ["dave", "eve"]
        assert item.days_waiting == 2

    def test_contributed_pr_shows_author(self):
        authored_beta = self.content[0].blocks[1]  # org/beta
        item = authored_beta.items[0]
        assert item.author == "bob"


class TestRegroupProject:
    """Tests for group_by='project' mode."""

    @pytest.fixture(autouse=True)
    def _setup(self):
        report = _make_group_by_report()
        self.content = regroup_content(report, "project")

    def test_top_level_groups(self):
        names = [rc.repo_name for rc in self.content]
        assert names == ["org/alpha", "org/beta"]

    def test_alpha_has_open_status_block(self):
        alpha = self.content[0]
        headings = [b.heading for b in alpha.blocks]
        assert "Open" in headings
        open_block = [b for b in alpha.blocks if b.heading == "Open"][0]
        # Should have 2 items: authored "Add login" + reviewed "Update docs"
        assert len(open_block.items) == 2
        titles = {i.title for i in open_block.items}
        assert "Add login" in titles
        assert "Update docs" in titles

    def test_beta_has_merged_and_waiting_blocks(self):
        beta = self.content[1]
        headings = [b.heading for b in beta.blocks]
        assert "Merged" in headings
        assert "Waiting for Review" in headings

    def test_status_order(self):
        beta = self.content[1]
        headings = [b.heading for b in beta.blocks]
        # Merged should come before Waiting for Review (per _STATUS_ORDER)
        assert headings.index("Merged") < headings.index("Waiting for Review")

    def test_item_fields_preserved(self):
        alpha = self.content[0]
        open_block = [b for b in alpha.blocks if b.heading == "Open"][0]
        login_item = [i for i in open_block.items if i.title == "Add login"][0]
        assert login_item.additions == 50
        assert login_item.deletions == 10
        assert login_item.numbers == [10]


class TestRegroupStatus:
    """Tests for group_by='status' mode."""

    @pytest.fixture(autouse=True)
    def _setup(self):
        report = _make_group_by_report()
        self.content = regroup_content(report, "status")

    def test_top_level_groups(self):
        names = [rc.repo_name for rc in self.content]
        assert names == ["Open", "Merged", "Waiting for Review"]

    def test_open_group_has_alpha_block(self):
        open_group = self.content[0]
        headings = [b.heading for b in open_group.blocks]
        assert "org/alpha" in headings
        alpha_block = [b for b in open_group.blocks if b.heading == "org/alpha"][0]
        assert len(alpha_block.items) == 2
        titles = {i.title for i in alpha_block.items}
        assert "Add login" in titles
        assert "Update docs" in titles

    def test_merged_group_has_beta_block(self):
        merged_group = self.content[1]
        headings = [b.heading for b in merged_group.blocks]
        assert "org/beta" in headings
        beta_block = [b for b in merged_group.blocks if b.heading == "org/beta"][0]
        assert len(beta_block.items) == 1
        assert beta_block.items[0].title == "Fix crash"

    def test_waiting_group_has_beta_block(self):
        waiting_group = self.content[2]
        headings = [b.heading for b in waiting_group.blocks]
        assert "org/beta" in headings

    def test_item_status_cleared(self):
        # In status mode, item.status should be empty since the parent is the status
        open_group = self.content[0]
        for block in open_group.blocks:
            for item in block.items:
                assert item.status == ""

        merged_group = self.content[1]
        for block in merged_group.blocks:
            for item in block.items:
                assert item.status == ""

    def test_waiting_items_retain_reviewers(self):
        waiting_group = self.content[2]
        beta_block = [b for b in waiting_group.blocks if b.heading == "org/beta"][0]
        item = beta_block.items[0]
        assert item.reviewers == ["dave", "eve"]
        assert item.days_waiting == 2


class TestRegroupEmpty:
    """Tests for edge cases: empty report, single category."""

    def test_empty_report_returns_empty(self):
        report = ReportData(
            user="alice",
            date_from="2026-02-10",
            date_to="2026-02-10",
            summary=SummaryStats(
                total_prs=0, repo_count=0, merged_count=0,
                open_count=0, themes=[], is_range=False,
            ),
        )
        for mode in ("project", "status", "contribution"):
            result = regroup_content(report, mode)
            assert result == [], f"Expected empty for mode={mode}"

    def test_single_category_only(self):
        report = ReportData(
            user="alice",
            date_from="2026-02-10",
            date_to="2026-02-10",
            authored_prs=[
                AuthoredPR(
                    repo="org/alpha", title="Add login", number=10,
                    status="Open", additions=50, deletions=10,
                    contributed=False, original_author=None,
                ),
            ],
            summary=SummaryStats(
                total_prs=1, repo_count=1, merged_count=0,
                open_count=1, themes=[], is_range=False,
            ),
        )

        # Contribution mode: only "Authored / Contributed" group
        content = regroup_content(report, "contribution")
        names = [rc.repo_name for rc in content]
        assert names == ["Authored / Contributed"]

        # Project mode: only "org/alpha" group
        content = regroup_content(report, "project")
        names = [rc.repo_name for rc in content]
        assert names == ["org/alpha"]

        # Status mode: only "Open" group
        content = regroup_content(report, "status")
        names = [rc.repo_name for rc in content]
        assert names == ["Open"]


# ===========================================================================
# 2. Markdown formatting tests
# ===========================================================================


class TestFormatMarkdownProject:
    """Markdown output with group_by='project'."""

    @pytest.fixture(autouse=True)
    def _setup(self):
        self.report = _make_regrouped_report("project")
        self.md = format_markdown(self.report, group_by="project")

    def test_header_has_backticks(self):
        assert "## `org/alpha`" in self.md

    def test_block_heading_no_backticks(self):
        assert "- **Open**" in self.md
        # Should NOT have backticks around status heading
        assert "**`Open`**" not in self.md

    def test_pr_link_uses_repo_name(self):
        assert "https://github.com/org/alpha/pull/10" in self.md


class TestFormatMarkdownStatus:
    """Markdown output with group_by='status'."""

    @pytest.fixture(autouse=True)
    def _setup(self):
        self.report = _make_regrouped_report("status")
        self.md = format_markdown(self.report, group_by="status")

    def test_header_no_backticks(self):
        assert "## Open" in self.md
        assert "## `Open`" not in self.md

    def test_block_heading_has_backticks(self):
        assert "- **`org/alpha`**" in self.md


class TestFormatMarkdownContribution:
    """Markdown output with group_by='contribution'."""

    @pytest.fixture(autouse=True)
    def _setup(self):
        self.report = _make_regrouped_report("contribution")
        self.md = format_markdown(self.report, group_by="contribution")

    def test_header_no_backticks(self):
        assert "## Authored / Contributed" in self.md
        assert "## `Authored" not in self.md

    def test_block_heading_has_backticks(self):
        assert "- **`org/alpha`**" in self.md

    def test_pr_link_uses_block_heading(self):
        # In contribution mode, block.heading is the project name
        # PR links should use the block heading (project) for the URL
        assert "https://github.com/org/alpha/pull/10" in self.md
        assert "https://github.com/org/beta/pull/20" in self.md


# ===========================================================================
# 3. Slides formatting tests
# ===========================================================================


class TestFormatSlidesGroupBy:
    """Slide deck generation respects group_by parameter."""

    @pytest.fixture(autouse=True)
    def _generate(self, tmp_path):
        self.tmp_path = tmp_path

    def test_project_mode_slide_titles(self):
        report = _make_regrouped_report("project")
        output_path = str(self.tmp_path / "project.pptx")
        format_slides(report, output_path, group_by="project")
        prs = Presentation(output_path)
        # Skip title (0) and summary (last), content slides in between
        content_titles = [
            prs.slides[i].shapes.title.text
            for i in range(1, len(prs.slides) - 1)
        ]
        assert "org/alpha" in content_titles
        assert "org/beta" in content_titles

    def test_status_mode_slide_titles(self):
        report = _make_regrouped_report("status")
        output_path = str(self.tmp_path / "status.pptx")
        format_slides(report, output_path, group_by="status")
        prs = Presentation(output_path)
        content_titles = [
            prs.slides[i].shapes.title.text
            for i in range(1, len(prs.slides) - 1)
        ]
        assert "Open" in content_titles
        assert "Merged" in content_titles
        assert "Waiting for Review" in content_titles

    def test_contribution_mode_slide_titles(self):
        report = _make_regrouped_report("contribution")
        output_path = str(self.tmp_path / "contribution.pptx")
        format_slides(report, output_path, group_by="contribution")
        prs = Presentation(output_path)
        content_titles = [
            prs.slides[i].shapes.title.text
            for i in range(1, len(prs.slides) - 1)
        ]
        assert "Authored / Contributed" in content_titles
        assert "Reviewed" in content_titles
        assert "Waiting for Review" in content_titles


# ===========================================================================
# 4. Slack formatting tests
# ===========================================================================


class TestFormatSlackGroupBy:
    """Slack Block Kit output respects group_by parameter."""

    def test_project_mode_repo_header(self):
        report = _make_regrouped_report("project")
        result = format_slack(report, group_by="project")
        texts = _all_section_texts(result["blocks"])
        combined = "\n".join(texts)
        assert "*`org/alpha`*" in combined

    def test_status_mode_repo_header(self):
        report = _make_regrouped_report("status")
        result = format_slack(report, group_by="status")
        texts = _all_section_texts(result["blocks"])
        combined = "\n".join(texts)
        # Status mode: no backticks around status names
        assert "*Open*" in combined
        assert "*`Open`*" not in combined

    def test_contribution_mode_repo_header(self):
        report = _make_regrouped_report("contribution")
        result = format_slack(report, group_by="contribution")
        texts = _all_section_texts(result["blocks"])
        combined = "\n".join(texts)
        assert "*Authored / Contributed*" in combined
        assert "*`Authored" not in combined

    def test_project_mode_block_heading(self):
        report = _make_regrouped_report("project")
        result = format_slack(report, group_by="project")
        texts = _all_section_texts(result["blocks"])
        combined = "\n".join(texts)
        # Block headings in project mode are statuses (L1 bullets), no backticks
        assert "\u2022 *Open*" in combined
        assert "*`Open`*" not in combined

    def test_status_mode_block_heading(self):
        report = _make_regrouped_report("status")
        result = format_slack(report, group_by="status")
        texts = _all_section_texts(result["blocks"])
        combined = "\n".join(texts)
        # Block headings in status mode are project names (L1 bullets), with backticks
        assert "\u2022 *`org/alpha`*" in combined


# ===========================================================================
# 5. CLI argument tests
# ===========================================================================


class TestCLIGroupBy:
    """CLI argument validation for --group-by."""

    def test_group_by_help_text(self):
        cmd = [sys.executable, "-m", "daily_report", "--help"]
        result = subprocess.run(
            cmd, capture_output=True, text=True, timeout=30,
            cwd=str(PROJECT_ROOT),
        )
        assert "group-by" in result.stdout

    def test_invalid_group_by_value(self):
        cmd = [
            sys.executable, "-m", "daily_report",
            "--group-by", "invalid_mode",
        ]
        result = subprocess.run(
            cmd, capture_output=True, text=True, timeout=30,
            cwd=str(PROJECT_ROOT),
        )
        assert result.returncode != 0
        assert "invalid choice" in result.stderr
